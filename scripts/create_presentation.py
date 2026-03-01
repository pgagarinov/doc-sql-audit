"""Generate a PowerPoint presentation for the AI procurement audit demo.

Usage:
    pixi run python scripts/create_presentation.py

Output: presentation.pptx in the project root.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from lxml import etree

# ---------------------------------------------------------------------------
# Design constants — harmonious green theme
# ---------------------------------------------------------------------------
# Gradient
GREEN_TOP = RGBColor(0xE8, 0xF5, 0xE9)   # light green
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# Text hierarchy
DARK_TEXT = RGBColor(0x1B, 0x5E, 0x20)    # dark green for titles
BODY_TEXT = RGBColor(0x37, 0x47, 0x4F)    # blue-gray body text
SUBTLE_TEXT = RGBColor(0x60, 0x7D, 0x8B)  # muted teal-gray

# Accent — forest green palette
ACCENT = RGBColor(0x2E, 0x7D, 0x32)       # primary accent (forest green)
ACCENT_LIGHT = RGBColor(0x66, 0xBB, 0x6A) # lighter green for charts
WARN_RED = RGBColor(0xE5, 0x73, 0x73)     # muted coral for violations

# Surfaces
CARD_BG = RGBColor(0xFF, 0xFF, 0xFF)
CARD_BORDER = RGBColor(0xA5, 0xD6, 0xA7)  # medium-light green

# Table
HEADER_BG = RGBColor(0x2E, 0x7D, 0x32)    # forest green header
ROW_EVEN = RGBColor(0xFF, 0xFF, 0xFF)
ROW_ODD = RGBColor(0xF1, 0xF8, 0xF1)      # very light green
TABLE_BORDER_HEX = "C8E6C9"

# Chat bubbles
BUBBLE_USER_BG = RGBColor(0x2E, 0x7D, 0x32)  # forest green
BUBBLE_AI_BG = RGBColor(0xFF, 0xFF, 0xFF)

# Layout
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
FONT_NAME = "Calibri"


# ---------------------------------------------------------------------------
# Helper functions
# ---------------------------------------------------------------------------

def _set_bg(slide):
    """Fill slide background with gradient from light green (top) to white (bottom)."""
    bg = slide.background._element  # <p:bg>
    # Remove existing fill settings
    for tag in (qn('p:bgPr'), qn('p:bgRef')):
        for child in bg.findall(tag):
            bg.remove(child)
    # Build gradient XML
    bgPr = etree.SubElement(bg, qn('p:bgPr'))
    gradFill = etree.SubElement(bgPr, qn('a:gradFill'))
    gsLst = etree.SubElement(gradFill, qn('a:gsLst'))
    # Stop 1: light green at top
    gs1 = etree.SubElement(gsLst, qn('a:gs'))
    gs1.set('pos', '0')
    clr1 = etree.SubElement(gs1, qn('a:srgbClr'))
    clr1.set('val', 'E8F5E9')
    # Stop 2: white at bottom
    gs2 = etree.SubElement(gsLst, qn('a:gs'))
    gs2.set('pos', '100000')
    clr2 = etree.SubElement(gs2, qn('a:srgbClr'))
    clr2.set('val', 'FFFFFF')
    # Linear direction: top to bottom
    lin = etree.SubElement(gradFill, qn('a:lin'))
    lin.set('ang', '5400000')  # 90 degrees
    lin.set('scaled', '1')
    etree.SubElement(bgPr, qn('a:effectLst'))


def _add_title(slide, text, *, top=Inches(0.4), left=Inches(0.7),
               width=Inches(12), height=Inches(0.9), size=Pt(36),
               color=None, bold=True, align=PP_ALIGN.LEFT):
    """Add a styled title textbox."""
    if color is None:
        color = DARK_TEXT
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = size
    p.font.color.rgb = color
    p.font.name = FONT_NAME
    p.font.bold = bold
    p.alignment = align
    return txBox


def _add_text(slide, text, *, top, left=Inches(0.7), width=Inches(11.9),
              height=Inches(0.6), size=Pt(20), color=None,
              bold=False, align=PP_ALIGN.LEFT):
    """Add a body text box."""
    if color is None:
        color = BODY_TEXT
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = size
    p.font.color.rgb = color
    p.font.name = FONT_NAME
    p.font.bold = bold
    p.alignment = align
    return txBox


def _add_bullet_list(slide, items, *, top, left=Inches(0.7),
                     width=Inches(11.9), height=None, size=Pt(20),
                     color=None, bullet_color=None):
    """Add a bulleted list of strings."""
    if color is None:
        color = BODY_TEXT
    if bullet_color is None:
        bullet_color = ACCENT
    if height is None:
        height = Inches(len(items) * 0.5 + 0.2)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = size
        p.font.color.rgb = color
        p.font.name = FONT_NAME
        p.space_after = Pt(6)
        p.level = 0
        pPr = p._pPr
        if pPr is None:
            pPr = p._p.get_or_add_pPr()
        for b in pPr.findall(qn("a:buNone")):
            pPr.remove(b)
        buChar = pPr.find(qn("a:buChar"))
        if buChar is None:
            buChar = etree.SubElement(pPr, qn("a:buChar"))
        buChar.set("char", "\u25cf")
        buClr = pPr.find(qn("a:buClr"))
        if buClr is None:
            buClr = etree.SubElement(pPr, qn("a:buClr"))
        else:
            buClr.clear()
        srgb = etree.SubElement(buClr, qn("a:srgbClr"))
        srgb.set("val", str(bullet_color))
    return txBox


def _add_table(slide, headers, rows, *,
               top, left=Inches(0.7), width=Inches(11.9),
               row_height=Inches(0.45), col_widths=None,
               header_size=Pt(16), cell_size=Pt(14)):
    """Add a styled table with light theme."""
    n_rows = len(rows) + 1
    n_cols = len(headers)
    height = row_height * n_rows

    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = table_shape.table
    tbl = table._tbl
    tblPr = tbl.tblPr
    tblPr.set("bandRow", "0")
    tblPr.set("firstRow", "0")
    tblPr.set("lastRow", "0")

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    def _style_cell(cell, text, bg, fg, font_size, bold=False, align=PP_ALIGN.LEFT):
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        p.text = str(text)
        p.font.size = font_size
        p.font.color.rgb = fg
        p.font.name = FONT_NAME
        p.font.bold = bold
        p.alignment = align
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        tcPr = cell._tc.get_or_add_tcPr()
        for old in tcPr.findall(qn("a:solidFill")):
            tcPr.remove(old)
        solidFill = etree.SubElement(tcPr, qn("a:solidFill"))
        srgb = etree.SubElement(solidFill, qn("a:srgbClr"))
        srgb.set("val", str(bg))
        for border_name in ("a:lnL", "a:lnR", "a:lnT", "a:lnB"):
            ln = tcPr.find(qn(border_name))
            if ln is None:
                ln = etree.SubElement(tcPr, qn(border_name))
            ln.set("w", "6350")
            sfill = ln.find(qn("a:solidFill"))
            if sfill is None:
                sfill = etree.SubElement(ln, qn("a:solidFill"))
            else:
                sfill.clear()
            sc = etree.SubElement(sfill, qn("a:srgbClr"))
            sc.set("val", TABLE_BORDER_HEX)

    # Header
    for j, h in enumerate(headers):
        _style_cell(table.cell(0, j), h, HEADER_BG, WHITE, header_size,
                     bold=True, align=PP_ALIGN.CENTER)
    # Data rows
    for i, row in enumerate(rows):
        bg = ROW_EVEN if i % 2 == 0 else ROW_ODD
        for j, val in enumerate(row):
            _style_cell(table.cell(i + 1, j), val, bg, DARK_TEXT, cell_size)

    return table_shape


def _add_accent_line(slide, top=Inches(1.3), left=Inches(0.7),
                     width=Inches(2.5)):
    """Add a short orange accent line under the title."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, Pt(4)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()


def _add_stat_card(slide, value, label, *,
                   left, top, width=Inches(3.2), height=Inches(2.0)):
    """Add a rounded card with a large number and label (white card, green border)."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    shape.line.color.rgb = CARD_BORDER
    shape.line.width = Pt(1.5)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_top = Inches(0.25)

    p = tf.paragraphs[0]
    p.text = value
    p.font.size = Pt(36)
    p.font.color.rgb = ACCENT
    p.font.name = FONT_NAME
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = label
    p2.font.size = Pt(16)
    p2.font.color.rgb = BODY_TEXT
    p2.font.name = FONT_NAME
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(8)


def _add_chat_bubble(slide, text, *, top, left, width, height,
                     bg_color, text_color=None, size=Pt(18),
                     label=None, label_color=None, border_color=None):
    """Add a chat-style bubble with optional label above."""
    if text_color is None:
        text_color = WHITE if bg_color == BUBBLE_USER_BG else DARK_TEXT
    if label:
        lbl_color = label_color or SUBTLE_TEXT
        _add_text(slide, label, top=top - Inches(0.3), left=left,
                  width=width, height=Inches(0.3), size=Pt(13),
                  color=lbl_color, bold=True)
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    shape.adjustments[0] = 0.06
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.25)
    tf.margin_right = Inches(0.25)
    tf.margin_top = Inches(0.15)
    tf.margin_bottom = Inches(0.15)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = size
    p.font.color.rgb = text_color
    p.font.name = FONT_NAME
    p.alignment = PP_ALIGN.LEFT
    return shape


def _add_arch_card(slide, title, desc, *, left, top,
                   width=Inches(2.6), height=Inches(2.0),
                   highlight=False):
    """Add an architecture card with title and description."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    shape.line.color.rgb = ACCENT if highlight else CARD_BORDER
    shape.line.width = Pt(2) if highlight else Pt(1.5)

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0.25)
    tf.margin_left = Inches(0.15)
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(22)
    p.font.color.rgb = ACCENT
    p.font.name = FONT_NAME
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(16)
    p2.font.color.rgb = BODY_TEXT
    p2.font.name = FONT_NAME
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(12)
    return shape


def _add_db_table_box(slide, name, columns, *, left, top,
                      width=Inches(3.5), height=Inches(2.8)):
    """Draw a database table box for ER diagram."""
    # Header bar (table name)
    hdr_h = Inches(0.5)
    hdr = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, hdr_h
    )
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = ACCENT
    hdr.line.color.rgb = ACCENT
    hdr.line.width = Pt(1)
    tf = hdr.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_top = Inches(0.05)
    p = tf.paragraphs[0]
    p.text = name
    p.font.size = Pt(18)
    p.font.color.rgb = WHITE
    p.font.name = FONT_NAME
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Body (columns)
    body_top = top + hdr_h - Pt(2)
    body_h = height - hdr_h + Pt(2)
    body = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, body_top, width, body_h
    )
    body.fill.solid()
    body.fill.fore_color.rgb = WHITE
    body.line.color.rgb = CARD_BORDER
    body.line.width = Pt(1.5)

    tf = body.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_top = Inches(0.12)
    for i, col_text in enumerate(columns):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.font.size = Pt(12)
        p.font.name = FONT_NAME
        p.space_before = Pt(2)
        p.space_after = Pt(1)
        # Bold for PK/FK markers
        if col_text.startswith("PK") or col_text.startswith("FK"):
            run_marker = p.add_run()
            marker, rest = col_text.split(" ", 1)
            run_marker.text = marker + " "
            run_marker.font.size = Pt(11)
            run_marker.font.color.rgb = ACCENT
            run_marker.font.name = FONT_NAME
            run_marker.font.bold = True
            run_field = p.add_run()
            run_field.text = rest
            run_field.font.size = Pt(12)
            run_field.font.color.rgb = DARK_TEXT
            run_field.font.name = FONT_NAME
            run_field.font.bold = col_text.startswith("PK")
        else:
            p.text = col_text
            p.font.color.rgb = BODY_TEXT
    return body


def _add_arrow_text(slide, *, top, left, label="", direction="right"):
    """Add an arrow with optional label between ER boxes."""
    arrow = "\u2192" if direction == "right" else "\u2190"
    w = Inches(0.8)
    # Arrow
    _add_text(slide, arrow, top=top, left=left,
              width=w, height=Inches(0.4),
              size=Pt(28), color=ACCENT, bold=True,
              align=PP_ALIGN.CENTER)
    # Label below arrow
    if label:
        _add_text(slide, label, top=top + Inches(0.3), left=left,
                  width=w, height=Inches(0.3),
                  size=Pt(9), color=SUBTLE_TEXT, bold=False,
                  align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def slide_01_title(prs):
    """Slide 1: Title."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)

    _add_title(slide, "ИИ-аудит процесса\nзаключения договоров",
               top=Inches(1.8), left=Inches(0.9), size=Pt(48),
               align=PP_ALIGN.LEFT)
    _add_accent_line(slide, top=Inches(3.4), left=Inches(0.9), width=Inches(3))

    _add_text(slide,
              "Автоматический поиск нарушений на естественном языке",
              top=Inches(3.7), left=Inches(0.9), size=Pt(24),
              color=BODY_TEXT)

    _add_text(slide,
              "100 000 договоров  \u00b7  500 поставщиков  \u00b7  ~17,5 млрд руб.",
              top=Inches(5.0), left=Inches(0.9), size=Pt(20),
              color=ACCENT, bold=True)


def slide_02_problem(prs):
    """Slide 2: Problem."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)

    _add_title(slide, "Почему ручной аудит не справляется")
    _add_accent_line(slide)

    items = [
        "100 000 договоров \u2014 невозможно проверить вручную",
        "18 форматов доп. соглашений без единой структуры",
        "Сложные связи между документами, протоколами, согласованиями",
        "Нужен аудитор + программист + время = дорого и медленно",
    ]
    _add_bullet_list(slide, items, top=Inches(1.8), size=Pt(22),
                     width=Inches(11.0))

    _add_stat_card(slide, "100 000", "договоров\nдля проверки",
                   left=Inches(0.7), top=Inches(4.5), width=Inches(3.5))
    _add_stat_card(slide, "18", "форматов\nдоп. соглашений",
                   left=Inches(4.8), top=Inches(4.5), width=Inches(3.5))
    _add_stat_card(slide, "3", "таблицы\nсо связями",
                   left=Inches(8.9), top=Inches(4.5), width=Inches(3.5))


def slide_03_solution(prs):
    """Slide 3: Solution — chat UX."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)

    _add_title(slide, "Как это работает")
    _add_accent_line(slide)

    # User bubble
    _add_chat_bubble(
        slide,
        "Найди группы нестандартных договоров свыше 10 000 руб.\n"
        "без согласования юриста",
        top=Inches(2.2), left=Inches(0.8), width=Inches(9.0), height=Inches(1.0),
        bg_color=BUBBLE_USER_BG, text_color=WHITE, size=Pt(20),
        label="\U0001F464  Аудитор"
    )

    # AI response bubble
    _add_chat_bubble(
        slide,
        "Найдено 21 204 договора в 7 063 группах с нарушениями.\n"
        "Общая сумма: 3,8 млрд руб. Вот первые результаты:",
        top=Inches(4.0), left=Inches(3.5), width=Inches(9.0), height=Inches(1.0),
        bg_color=BUBBLE_AI_BG, text_color=DARK_TEXT, size=Pt(20),
        label="\U0001F916  Система",
        border_color=CARD_BORDER
    )

    _add_text(slide,
              "Без SQL.  Без программирования.  Без ожидания.",
              top=Inches(5.8), left=Inches(0.7), width=Inches(11.9),
              size=Pt(28), color=ACCENT, bold=True, align=PP_ALIGN.CENTER)


def slide_04_demo_table(prs):
    """Slide 4: Demo query -> table with subject column."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)

    _add_title(slide, "Демо: запрос \u2192 результат")
    _add_accent_line(slide)

    # User query bubble
    _add_chat_bubble(
        slide,
        "Покажи поставщиков с признаками дробления закупок",
        top=Inches(1.65), left=Inches(0.8), width=Inches(8.0), height=Inches(0.55),
        bg_color=BUBBLE_USER_BG, text_color=WHITE, size=Pt(17),
        label="\U0001F464  Запрос аудитора"
    )

    # 6 columns: №, Поставщик, Предмет, Номера договоров, Сумма, Кол-во
    total_w = Inches(12.0)
    col_widths = [
        Inches(0.5),   # №
        Inches(1.7),   # Поставщик
        Inches(2.5),   # Предмет
        Inches(4.5),   # Номера договоров
        Inches(1.6),   # Сумма
        Inches(1.2),   # Кол-во
    ]

    headers = ["\u2116", "Поставщик", "Предмет", "Номера договоров",
               "Сумма, руб.", "Кол-во"]
    rows = [
        ["1", "ООО \u00abЛогист\u00bb",
         "Поставка канцтоваров,\nпоставка бумаги",
         "Договор 128, Договор 4501,\nДоговор 9823, дс к дог. 128",
         "2 450 300", "4"],
        ["2", "ЗАО \u00abСнаб\u00bb",
         "Поставка компьютерной\nтехники",
         "Договор 512, Договор 7744,\nдоп. соглашение к дог. 512",
         "1 870 000", "3"],
        ["3", "АО \u00abТехно\u00bb",
         "Поставка серверного\nоборудования",
         "Договор 3901, Договор 3902,\nДС \u21161 к дог. 3901",
         "1 340 500", "3"],
        ["4", "ИП Иванов",
         "Услуги по ремонту\nавтотранспорта",
         "Договор 221, Договор 6680",
         "980 200", "2"],
        ["5", "ООО \u00abСервис\u00bb",
         "Оказание клининговых\nуслуг",
         "Договор 1055, Договор 1056,\nДоговор 1057, дс к 1055",
         "875 100", "4"],
        ["6", "ЗАО \u00abАльфа\u00bb",
         "Поставка офисной\nмебели",
         "Договор 8801, Договор 8802",
         "720 400", "2"],
    ]

    _add_table(slide, headers, rows,
               top=Inches(2.65), left=Inches(0.65), width=total_w,
               row_height=Inches(0.6),
               col_widths=col_widths,
               header_size=Pt(13), cell_size=Pt(11))

    _add_text(slide,
              "Система сама находит связанные договоры, "
              "группирует по поставщикам и сортирует по сумме",
              top=Inches(6.95), size=Pt(15), color=SUBTLE_TEXT)


def slide_05_check1(prs):
    """Slide 5: Check 1 — non-standard without lawyer."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)

    _add_title(slide, "Проверка 1: Нестандартные без юриста")
    _add_accent_line(slide)

    _add_chat_bubble(
        slide,
        "Найди группы нестандартных договоров без согласования юриста",
        top=Inches(1.65), left=Inches(0.8), width=Inches(9.5), height=Inches(0.55),
        bg_color=BUBBLE_USER_BG, text_color=WHITE, size=Pt(17),
        label="\U0001F464  Запрос аудитора"
    )

    _add_stat_card(slide, "21 204", "договора в группах\nс нарушениями",
                   left=Inches(0.7), top=Inches(2.8), width=Inches(3.5),
                   height=Inches(1.7))

    _add_text(slide, "Что проверяет система:",
              top=Inches(2.8), left=Inches(4.8), size=Pt(15),
              color=ACCENT, bold=True, width=Inches(7.5))

    items = [
        "Группирует ДС с основным договором автоматически",
        "Сумма группы > 10 000 руб., тип \u00abнестандартный\u00bb",
        "Нет согласования юриста в хотя бы одном документе",
        "Протокол не разрешает работу без юриста",
    ]
    _add_bullet_list(slide, items, top=Inches(3.2), left=Inches(4.8),
                     width=Inches(7.5), size=Pt(15))

    headers = ["Договор", "Поставщик", "Сумма", "Тип"]
    sample_rows = [
        ["Договор 1", "ЗАО \u00abФомина, Новикова\nи Семенова\u00bb",
         "202 715", "нестандартный"],
        ["доп. соглашение\nк дог. 1", "ЗАО \u00abФомина, Новикова\nи Семенова\u00bb",
         "121 218", "нестандартный"],
        ["Договор 2", "АО \u00abЛазарева, Воронов\nи Шарова\u00bb",
         "44 779", "стандартный"],
    ]
    _add_table(slide, headers, sample_rows,
               top=Inches(5.3), left=Inches(0.7), width=Inches(11.9),
               row_height=Inches(0.55),
               col_widths=[Inches(2.8), Inches(4.0), Inches(2.0), Inches(3.1)],
               header_size=Pt(13), cell_size=Pt(12))


def slide_06_check2(prs):
    """Slide 6: Check 2 — splitting detection (pairs)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)

    _add_title(slide, "Проверка 2: Дробление закупок")
    _add_accent_line(slide)

    _add_chat_bubble(
        slide,
        "Найди пары договоров одного поставщика со схожими предметами",
        top=Inches(1.65), left=Inches(0.8), width=Inches(9.5), height=Inches(0.55),
        bg_color=BUBBLE_USER_BG, text_color=WHITE, size=Pt(17),
        label="\U0001F464  Запрос аудитора"
    )

    _add_stat_card(slide, "1 019 115", "подозрительных пар",
                   left=Inches(0.7), top=Inches(2.8), width=Inches(3.8),
                   height=Inches(1.7))

    _add_text(slide, "Как работает поиск:",
              top=Inches(2.8), left=Inches(5.2), size=Pt(15),
              color=ACCENT, bold=True, width=Inches(7.0))

    items = [
        "Сравнивает предметы договоров одного поставщика",
        "Нечёткий поиск: \u00abпоставка канцтоваров\u00bb и\n"
        "\u00abпоставка канцелярских принадлежностей\u00bb совпадут",
        "Аудитор просто получает готовый результат",
    ]
    _add_bullet_list(slide, items, top=Inches(3.2), left=Inches(5.2),
                     width=Inches(7.0), size=Pt(15))

    _add_text(slide, "Пример найденной пары:",
              top=Inches(5.2), left=Inches(0.7), size=Pt(15),
              color=ACCENT, bold=True)

    headers = ["Предмет договора 1", "Предмет договора 2", "Поставщик", "Схожесть"]
    pair_rows = [
        ["Поставка канцелярских\nтоваров",
         "Поставка канцелярских\nпринадлежностей",
         "ООО \u00abОфис-М\u00bb", "78%"],
    ]
    _add_table(slide, headers, pair_rows,
               top=Inches(5.6), left=Inches(0.7), width=Inches(11.9),
               row_height=Inches(0.6),
               col_widths=[Inches(3.5), Inches(3.5), Inches(2.5), Inches(2.4)],
               header_size=Pt(13), cell_size=Pt(12))


def slide_07_check3(prs):
    """Slide 7: Check 3 — systematic splitting."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)

    _add_title(slide, "Проверка 3: Систематическое дробление")
    _add_accent_line(slide)

    _add_chat_bubble(
        slide,
        "Сгруппируй подозрительные договоры по поставщикам и предметам",
        top=Inches(1.65), left=Inches(0.8), width=Inches(9.5), height=Inches(0.55),
        bg_color=BUBBLE_USER_BG, text_color=WHITE, size=Pt(17),
        label="\U0001F464  Запрос аудитора"
    )

    _add_stat_card(slide, "19 801", "группа с признаками\nдробления",
                   left=Inches(0.7), top=Inches(3.0), width=Inches(3.5))
    _add_stat_card(slide, "~17,5 млрд", "общая сумма\nдоговоров, руб.",
                   left=Inches(4.8), top=Inches(3.0), width=Inches(3.5))
    _add_stat_card(slide, "500", "поставщиков\nв анализе",
                   left=Inches(8.9), top=Inches(3.0), width=Inches(3.5))

    _add_text(slide,
              "Система объединяет пары в группы: один поставщик + "
              "схожие предметы \u2192 единая группа с общей суммой",
              top=Inches(5.4), left=Inches(0.7), width=Inches(11.9),
              size=Pt(18), color=BODY_TEXT)

    _add_text(slide,
              "Аудитор видит картину целиком: кто, что, на какую сумму",
              top=Inches(6.1), left=Inches(0.7), width=Inches(11.9),
              size=Pt(20), color=ACCENT, bold=True)


def slide_08_results(prs):
    """Slide 8: Results — summary table + 3D pie chart."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)

    _add_title(slide, "Результаты: три проверки за минуты")
    _add_accent_line(slide)

    # Summary table (left side)
    headers = ["Проверка", "Что искали", "Результат"]
    rows = [
        ["1", "Нестандартные договоры\nбез согласования юриста", "21 204\nдоговора"],
        ["2", "Пары со схожими\nпредметами у поставщика", "1 019 115\nпар"],
        ["3", "Группы подозрительных\nдоговоров по поставщикам", "19 801\nгруппа"],
    ]
    _add_table(slide, headers, rows,
               top=Inches(1.6), left=Inches(0.5), width=Inches(7.2),
               row_height=Inches(0.85),
               col_widths=[Inches(1.2), Inches(3.5), Inches(2.5)],
               header_size=Pt(16), cell_size=Pt(14))

    # 3D pie chart (right side) — violation ratio
    # Create as regular PIE then patch XML to 3D
    chart_data = CategoryChartData()
    chart_data.categories = [
        'С нарушениями (21 204)',
        'Без нарушений (78 796)',
    ]
    chart_data.add_series('Договоры', (21204, 78796))

    chart_left = Inches(8.0)
    chart_top = Inches(1.5)
    chart_w = Inches(4.8)
    chart_h = Inches(3.8)

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE, chart_left, chart_top, chart_w, chart_h,
        chart_data
    )
    chart = chart_frame.chart
    chart.has_title = False
    chart.has_legend = True

    # Color segments and labels BEFORE patching to 3D
    plot = chart.plots[0]
    series = plot.series[0]
    pt0 = series.points[0]
    pt0.format.fill.solid()
    pt0.format.fill.fore_color.rgb = WARN_RED
    pt1 = series.points[1]
    pt1.format.fill.solid()
    pt1.format.fill.fore_color.rgb = ACCENT_LIGHT

    plot.has_data_labels = True
    data_labels = plot.data_labels
    data_labels.show_percentage = True
    data_labels.show_category_name = True
    data_labels.show_value = False
    data_labels.font.size = Pt(11)
    data_labels.font.color.rgb = DARK_TEXT

    # Patch XML: convert pieChart -> pie3DChart and add view3D
    chart_space = chart._element  # <c:chartSpace>
    chart_elem = chart_space.find(qn('c:chart'))
    plot_area = chart_elem.find(qn('c:plotArea'))
    pie_chart = plot_area.find(qn('c:pieChart'))
    if pie_chart is not None:
        pie_chart.tag = qn('c:pie3DChart')
    view3D = etree.SubElement(chart_elem, qn('c:view3D'))
    etree.SubElement(view3D, qn('c:rotX')).set('val', '30')
    etree.SubElement(view3D, qn('c:rotY')).set('val', '0')
    etree.SubElement(view3D, qn('c:perspective')).set('val', '30')
    chart_elem.insert(0, view3D)

    # Chart title above
    _add_text(slide, "Доля договоров с нарушениями",
              top=Inches(1.3), left=Inches(8.0), width=Inches(4.8),
              size=Pt(14), color=ACCENT, bold=True, align=PP_ALIGN.CENTER)

    # Bottom message
    _add_text(slide,
              "Ручная проверка заняла бы недели. Система справляется за минуты.",
              top=Inches(5.8), left=Inches(0.7), width=Inches(11.9),
              size=Pt(22), color=ACCENT, bold=True, align=PP_ALIGN.CENTER)


def slide_09_advantages(prs):
    """Slide 9: Advantages for auditor."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)

    _add_title(slide, "Преимущества для аудитора")
    _add_accent_line(slide)

    advantages = [
        ("Простота",
         "Запрос на русском языке, ответ в таблице. "
         "Не нужно знать SQL или структуру базы данных."),
        ("Скорость",
         "100 000+ договоров анализируются за секунды. "
         "Новая проверка \u2014 новый вопрос, не новый проект."),
        ("Гибкость",
         "Любая новая проверка без программиста. "
         "Аудитор сам формулирует, что искать."),
        ("Объяснимость",
         "Каждый результат можно проверить и перепроверить. "
         "Система прозрачна."),
        ("Масштабируемость",
         "Миллионы записей \u2014 без проблем. "
         "Работает с любым объёмом данных."),
    ]

    top = Inches(1.7)
    for i, (title, desc) in enumerate(advantages):
        y = top + Inches(i * 1.05)
        _add_text(slide, title, top=y, left=Inches(0.7),
                  width=Inches(11.9), size=Pt(22), color=ACCENT, bold=True)
        _add_text(slide, desc, top=y + Inches(0.35), left=Inches(0.7),
                  width=Inches(11.9), size=Pt(17), color=BODY_TEXT)


def slide_10_architecture(prs):
    """Slide 10: Architecture — for the technically curious."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)

    _add_title(slide, "Как это устроено")
    _add_accent_line(slide)

    _add_text(slide, "Для технически любопытных",
              top=Inches(1.5), size=Pt(16), color=SUBTLE_TEXT)

    cards = [
        ("Аудитор", "Задаёт вопрос\nна русском языке"),
        ("Claude", "Изучает схему БД,\nгенерирует запросы"),
        ("MCP-сервер", "Безопасный доступ\nк базе данных"),
        ("PostgreSQL", "100K договоров,\nнечёткий поиск"),
    ]
    box_w = Inches(2.6)
    box_h = Inches(2.0)
    gap = Inches(0.45)
    start_left = Inches(0.7)
    card_top = Inches(2.3)

    for i, (title, desc) in enumerate(cards):
        left = start_left + i * (box_w + gap)
        _add_arch_card(slide, title, desc, left=left, top=card_top,
                       width=box_w, height=box_h, highlight=(i == 1))
        if i < len(cards) - 1:
            arrow_left = left + box_w + Inches(0.05)
            arrow_top = card_top + box_h / 2 - Inches(0.15)
            _add_text(slide, "\u2192", top=arrow_top, left=arrow_left,
                      width=Inches(0.35), height=Inches(0.4),
                      size=Pt(28), color=ACCENT, bold=True,
                      align=PP_ALIGN.CENTER)

    _add_text(slide,
              "Claude автоматически изучает структуру базы данных, "
              "подбирает нужные таблицы и связи, формирует запрос "
              "и возвращает аудитору понятный результат.",
              top=Inches(4.8), left=Inches(0.7), width=Inches(11.9),
              size=Pt(18), color=BODY_TEXT)

    _add_text(slide,
              "Аудитор работает с интерфейсом чата. "
              "Вся техническая сложность скрыта.",
              top=Inches(5.5), left=Inches(0.7), width=Inches(11.9),
              size=Pt(20), color=ACCENT, bold=True)


def slide_11_tech_details(prs):
    """Slide 11: Technical details with ER diagram."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)

    _add_title(slide, "Технические детали")
    _add_accent_line(slide)

    _add_text(slide, "Структура данных",
              top=Inches(1.5), size=Pt(18), color=ACCENT, bold=True)

    # ER diagram: 3 table boxes with arrows
    box_w = Inches(3.3)
    box_h = Inches(2.5)
    box_top = Inches(1.9)

    # protocols (left)
    _add_db_table_box(slide, "protocols", [
        "PK protocol_id",
        "   protocol_number",
        "   protocol_date",
        "   approved_amount",
        "   protocol_text",
    ], left=Inches(0.5), top=box_top, width=box_w, height=box_h)

    # contracts (center)
    _add_db_table_box(slide, "contracts", [
        "PK contract_id",
        "   contract_number",
        "   amount",
        "   contract_type",
        "   supplier",
        "   subject",
        "FK protocol_id",
    ], left=Inches(4.8), top=box_top, width=box_w, height=box_h)

    # approvals (right)
    _add_db_table_box(slide, "approvals", [
        "PK approval_id",
        "FK contract_id",
        "   fin_director",
        "   lawyer",
        "   security",
        "   procurement_head",
    ], left=Inches(9.1), top=box_top, width=box_w, height=box_h)

    # Arrows between tables
    arrow_y = box_top + box_h / 2 - Inches(0.15)
    # contracts -> protocols
    _add_arrow_text(slide, top=arrow_y, left=Inches(3.85),
                    label="protocol_id", direction="left")
    # approvals -> contracts
    _add_arrow_text(slide, top=arrow_y, left=Inches(8.25),
                    label="contract_id", direction="left")

    # Key technologies below
    _add_text(slide, "Ключевые технологии",
              top=Inches(4.7), size=Pt(18), color=ACCENT, bold=True)

    items = [
        "pg_trgm \u2014 нечёткий поиск по предмету договора (similarity > 0.3)",
        "Группировка ДС: извлечение номера из 18 текстовых форматов",
        "GIN-индекс для ускорения триграммных запросов",
    ]
    _add_bullet_list(slide, items, top=Inches(5.1), size=Pt(16))

    _add_text(slide,
              "Настройка: 1 день   \u00b7   Поддержка: минимальная   \u00b7   "
              "Масштабирование: автоматическое",
              top=Inches(6.5), left=Inches(0.7), width=Inches(11.9),
              size=Pt(18), color=ACCENT, bold=True, align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_01_title(prs)
    slide_02_problem(prs)
    slide_03_solution(prs)
    slide_04_demo_table(prs)
    slide_05_check1(prs)
    slide_06_check2(prs)
    slide_07_check3(prs)
    slide_08_results(prs)
    slide_09_advantages(prs)
    slide_10_architecture(prs)
    slide_11_tech_details(prs)

    out = "presentation.pptx"
    prs.save(out)
    print(f"Saved \u2192 {out}")


if __name__ == "__main__":
    main()
